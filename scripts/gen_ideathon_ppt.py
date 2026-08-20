"""Generate AI Ideathon 2026 finals PowerPoint (4 min + 3 Q&A).

Reuse palette/helpers from the Round One deck; content follows
lessons/0006-ideathon-four-minute.html + SAFE numbers only.
"""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "lessons" / "0006-ideathon-final-pitch.pptx"

BG = RGBColor(0x0C, 0x06, 0x08)
INK = RGBColor(0xF5, 0xF5, 0xF4)
MUTED = RGBColor(0x8E, 0x8E, 0x96)
ACCENT = RGBColor(0xF2, 0x61, 0x3C)
BAD = RGBColor(0xFC, 0xA5, 0xA5)
OK = RGBColor(0x86, 0xEF, 0xAC)
CARD = RGBColor(0x16, 0x0A, 0x0E)
RULE = RGBColor(0x3F, 0x3F, 0x46)


def _blank_slide(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG
    return slide


def _box(slide, left, top, width, height, fill=CARD):
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = RULE
    shape.line.width = Pt(0.75)
    return shape


def _text(slide, left, top, width, height, text, size=14, bold=False, color=INK, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = align
    return tf


def _bullets(tf, items, size=14, color=INK):
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.level = 0
        p.space_after = Pt(8)


def _title_bar(slide, title: str, subtitle: str = ""):
    _text(slide, Inches(0.55), Inches(0.32), Inches(9), Inches(0.5), title, size=26, bold=True)
    bar = slide.shapes.add_shape(1, Inches(0.55), Inches(0.88), Inches(1.15), Inches(0.055))
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()
    if subtitle:
        _text(slide, Inches(0.55), Inches(0.98), Inches(9), Inches(0.3), subtitle, size=12, color=MUTED)


def slide_title(prs):
    slide = _blank_slide(prs)
    _text(
        slide,
        Inches(0.7),
        Inches(1.35),
        Inches(8.6),
        Inches(0.35),
        "AI Ideathon 2026  ·  Finals  ·  4 min + 3 Q&A",
        size=12,
        color=MUTED,
        align=PP_ALIGN.CENTER,
    )
    _text(slide, Inches(0.7), Inches(1.95), Inches(8.6), Inches(0.85), "aiand/auto", size=48, bold=True, align=PP_ALIGN.CENTER)
    _text(
        slide,
        Inches(1.1),
        Inches(2.9),
        Inches(7.8),
        Inches(0.7),
        "Cheapest aiand model per agent hop\nthat still clears your quality bar",
        size=18,
        color=MUTED,
        align=PP_ALIGN.CENTER,
    )
    _text(
        slide,
        Inches(0.7),
        Inches(4.15),
        Inches(8.6),
        Inches(0.35),
        "model: router/auto   ·   OpenAI-compatible gateway",
        size=14,
        color=ACCENT,
        align=PP_ALIGN.CENTER,
    )
    _text(
        slide,
        Inches(0.7),
        Inches(4.7),
        Inches(8.6),
        Inches(0.35),
        "Drop-in for OpenCode · Cursor · Claude Code · any OpenAI client",
        size=12,
        color=MUTED,
        align=PP_ALIGN.CENTER,
    )


def slide_problem(prs):
    slide = _blank_slide(prs)
    _title_bar(slide, "The problem", "~0:00–0:45  ·  Problem fit 25%")

    left = _box(slide, Inches(0.55), Inches(1.45), Inches(4.35), Inches(2.0))
    left.line.color.rgb = BAD
    tf = _text(slide, Inches(0.7), Inches(1.55), Inches(4.05), Inches(1.85), "", size=14)
    tf.paragraphs[0].text = "Premium-only"
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.size = Pt(16)
    tf.paragraphs[0].font.color.rgb = BAD
    _bullets(tf, ["Always K3 / frontier", "Burns $ on trivial greps", "Budget dies mid-session"], size=14, color=INK)

    right = _box(slide, Inches(5.1), Inches(1.45), Inches(4.35), Inches(2.0))
    right.line.color.rgb = BAD
    tf2 = _text(slide, Inches(5.25), Inches(1.55), Inches(4.05), Inches(1.85), "", size=14)
    tf2.paragraphs[0].text = "Cheap-only"
    tf2.paragraphs[0].font.bold = True
    tf2.paragraphs[0].font.size = Pt(16)
    tf2.paragraphs[0].font.color.rgb = BAD
    _bullets(tf2, ["Always Flash", "Fails hard refactors / debug", "Humans re-run broken loops"], size=14, color=INK)

    _text(slide, Inches(0.55), Inches(3.7), Inches(9), Inches(0.35), "Who feels it?", size=16, bold=True)
    tf3 = _text(slide, Inches(0.55), Inches(4.1), Inches(8.9), Inches(1.2), "", size=15)
    _bullets(
        tf3,
        [
            "Teams on OpenCode / Cursor / Claude Code × multi-model catalogs",
            "Agents make dozens of hops per session — not one call",
            "Impact: $ per session + fewer wasted human loops",
        ],
        size=15,
    )


def slide_solution(prs):
    slide = _blank_slide(prs)
    _title_bar(slide, "Solution & why AI", "~0:45–1:30  ·  Innovation 20%")

    tf = _text(slide, Inches(0.55), Inches(1.4), Inches(8.9), Inches(1.5), "", size=16)
    _bullets(
        tf,
        [
            "aiand/auto — OpenAI gateway · model: router/auto",
            "Each hop: filter ineligible → cheapest that clears quality bar",
            "Not static if phase=edit → Model X",
        ],
        size=16,
    )

    callout = _box(slide, Inches(0.55), Inches(3.15), Inches(8.9), Inches(2.0))
    callout.line.color.rgb = ACCENT
    _text(slide, Inches(0.7), Inches(3.25), Inches(8.6), Inches(0.35), "Purposeful AI", size=16, bold=True, color=ACCENT)
    tf2 = _text(slide, Inches(0.7), Inches(3.65), Inches(8.6), Inches(1.4), "", size=15)
    _bullets(
        tf2,
        [
            "Offline: teacher labels → gold on real aiand → fit Scorer",
            "Live: tiny student in ms — no LLM on the routing hop",
            "Better picks without taxing latency or cost of the route itself",
        ],
        size=15,
    )


def slide_demo(prs):
    slide = _blank_slide(prs)
    _title_bar(slide, "Live demo — one journey", "~1:30–2:45  ·  Demo 15%  ·  checklist only")

    steps = [
        ("1", "User need", "Cheap hop that still works"),
        ("2", "Action", "Playground (or OpenCode) → router/auto + phase"),
        ("3", "Show AI", "Chosen model · reason / score · path"),
        ("4", "Proof", "Dashboard: mix + savings vs expensive baseline"),
    ]
    y = 1.4
    for num, label, detail in steps:
        _text(slide, Inches(0.55), Inches(y), Inches(0.45), Inches(0.4), num, size=22, bold=True, color=ACCENT)
        _text(slide, Inches(1.1), Inches(y), Inches(2.0), Inches(0.4), label, size=16, bold=True)
        _text(slide, Inches(3.2), Inches(y), Inches(6.2), Inches(0.4), detail, size=16, color=MUTED)
        y += 0.55

    tip = _box(slide, Inches(0.55), Inches(3.7), Inches(8.9), Inches(1.4))
    tip.line.color.rgb = OK
    _text(slide, Inches(0.7), Inches(3.85), Inches(8.6), Inches(0.3), "Pre-open tabs · skip login / keys / menus", size=14, bold=True, color=OK)
    tf = _text(slide, Inches(0.7), Inches(4.25), Inches(8.6), Inches(0.75), "", size=14)
    _bullets(
        tf,
        [
            "Tabs: playground + /routers/auto",
            "Optional 5s landing simulator = vision/preview only — then live hop",
        ],
        size=14,
        color=MUTED,
    )


def slide_algorithm(prs):
    slide = _blank_slide(prs)
    _title_bar(slide, "How routing works", "~2:45–3:10  ·  Technical 20%")

    # Two paths as compact boxes
    _text(slide, Inches(0.55), Inches(1.3), Inches(4.4), Inches(0.28), "Rules path (router.py)", size=13, bold=True, color=ACCENT)
    tf = _text(slide, Inches(0.55), Inches(1.6), Inches(4.4), Inches(1.2), "", size=13)
    _bullets(
        tf,
        [
            "Hard constraints → phase bar",
            "Then Pioneer / cheap / max",
        ],
        size=13,
    )

    _text(slide, Inches(5.1), Inches(1.3), Inches(4.4), Inches(0.28), "Trained Scorer (scorer.py)", size=13, bold=True, color=ACCENT)
    tf2 = _text(slide, Inches(5.1), Inches(1.6), Inches(4.4), Inches(1.2), "", size=13)
    _bullets(
        tf2,
        [
            "Complexity bin → P(success)",
            "Cheapest-above-bar",
        ],
        size=13,
    )

    _text(
        slide,
        Inches(0.55),
        Inches(2.85),
        Inches(8.9),
        Inches(0.45),
        "Pioneer ≈ 0.40·P(success) + 0.20·capability + tools/latency/health − cost",
        size=12,
        color=MUTED,
    )

    _text(slide, Inches(0.55), Inches(3.35), Inches(9), Inches(0.28), "Offline pipeline (train.py)", size=13, bold=True)
    pipeline = [("Teacher", "labels"), ("Gold", "real runs"), ("Fit", "calibrate"), ("Shadow", "compare"), ("Flip", "gated")]
    x = 0.55
    for label, sub in pipeline:
        _box(slide, Inches(x), Inches(3.7), Inches(1.6), Inches(0.7))
        _text(slide, Inches(x + 0.05), Inches(3.75), Inches(1.5), Inches(0.3), label, size=12, bold=True, align=PP_ALIGN.CENTER)
        _text(slide, Inches(x + 0.05), Inches(4.05), Inches(1.5), Inches(0.25), sub, size=11, color=MUTED, align=PP_ALIGN.CENTER)
        x += 1.8

    _text(
        slide,
        Inches(0.55),
        Inches(4.65),
        Inches(8.9),
        Inches(0.4),
        "Hot path: in-process Scorer · no teacher LLM · target <10 ms",
        size=14,
        color=OK,
    )


def slide_metrics(prs):
    slide = _blank_slide(prs)
    _title_bar(slide, "How good is it — SAFE only", "Do not invent global savings %")

    rows = [
        ("Catalog", "9 aiand models"),
        ("Train queries", "2000"),
        ("scorer.json", "n_gold=1600 · n_silver=1000 · n_cal=800"),
        ("Artifact label", "not_spec_floors  —  not SWE-bench Verified"),
        ("Process spend", "≈ $8.16 (soft spend.txt)"),
        ("Tests", "213 pytest · CI FakeProvider never spends"),
        ("Savings log", "vs most_expensive_eligible per hop"),
    ]
    y = 1.35
    for label, val in rows:
        _text(slide, Inches(0.55), Inches(y), Inches(2.4), Inches(0.32), label, size=13, bold=True, color=ACCENT)
        _text(slide, Inches(3.0), Inches(y), Inches(6.4), Inches(0.32), val, size=13)
        y += 0.38

    _text(
        slide,
        Inches(0.55),
        Inches(4.2),
        Inches(8.9),
        Inches(0.9),
        "Honesty: rules are solid live; trained flips after operator gate.\n"
        "Landing 62% / 40–60% / hero % = vision/preview — not Verified wins.",
        size=13,
        color=MUTED,
    )


def slide_architecture(prs):
    slide = _blank_slide(prs)
    _title_bar(slide, "Architecture", "Say the boxes in ~15 seconds")

    boxes = [
        ("Agent", "OpenCode / Cursor"),
        ("FastAPI", "/v1/chat/completions"),
        ("Rules+Scorer", "router + scorer"),
        ("aiand", "9 models"),
        ("Log", "requests.jsonl"),
    ]
    x = 0.35
    for i, (title, sub) in enumerate(boxes):
        _box(slide, Inches(x), Inches(1.7), Inches(1.7), Inches(1.0))
        _text(slide, Inches(x + 0.05), Inches(1.8), Inches(1.6), Inches(0.35), title, size=13, bold=True, align=PP_ALIGN.CENTER)
        _text(slide, Inches(x + 0.05), Inches(2.2), Inches(1.6), Inches(0.4), sub, size=10, color=MUTED, align=PP_ALIGN.CENTER)
        if i < len(boxes) - 1:
            _text(slide, Inches(x + 1.7), Inches(1.95), Inches(0.25), Inches(0.4), "→", size=16, color=ACCENT, align=PP_ALIGN.CENTER)
        x += 1.95

    tf = _text(slide, Inches(0.55), Inches(3.2), Inches(8.9), Inches(1.8), "", size=15)
    _bullets(
        tf,
        [
            "Backend: FastAPI OpenAI-compatible proxy",
            "Frontend: Next.js landing · dashboard · playground",
            "Data: YAML catalog · scorer.json · JSONL telemetry",
            "Feasibility proof: 213 tests · no spend in CI",
        ],
        size=15,
    )


def slide_close(prs):
    slide = _blank_slide(prs)
    _title_bar(slide, "Potential & close", "~3:25–4:00  ·  Product potential 20%")

    tf = _text(slide, Inches(0.55), Inches(1.4), Inches(8.9), Inches(1.8), "", size=15)
    _bullets(
        tf,
        [
            "Users: teams & agent platforms already on OpenAI APIs",
            "Adoption = baseURL change — no harness rewrite",
            "Scale: add models in YAML → retrain offline → shadow → flip",
            "Next: promote after gates · SWE-bench Verified bar · SaaS / aiand",
        ],
        size=15,
    )

    close = _box(slide, Inches(0.55), Inches(3.5), Inches(8.9), Inches(1.5))
    close.line.color.rgb = ACCENT
    _text(slide, Inches(0.7), Inches(3.65), Inches(8.6), Inches(0.3), "Closing line", size=14, bold=True, color=ACCENT)
    _text(
        slide,
        Inches(0.7),
        Inches(4.05),
        Inches(8.6),
        Inches(0.85),
        "Agents will keep multiplying model calls. aiand/auto makes every hop\n"
        "cost-aware without rewriting the agent — and we can prove each pick in the log.",
        size=14,
    )


def slide_backup_qa1(prs):
    slide = _blank_slide(prs)
    _title_bar(slide, "BACKUP · Q&A — how / functional / vs", "≤25s answers")

    pairs = [
        ("How does the AI work?", "Teacher offline → gold real runs → fit · live scores in ms · no teacher call"),
        ("Is it functional?", "Yes — live rules · playground + dashboard · Scorer from JSON · 213 tests"),
        ("vs rule routers / OpenRouter?", "Per-request P(success) + cheapest-above-bar · telemetry · shadow before flip"),
    ]
    y = 1.4
    for q, a in pairs:
        _text(slide, Inches(0.55), Inches(y), Inches(8.9), Inches(0.3), q, size=14, bold=True, color=ACCENT)
        _text(slide, Inches(0.55), Inches(y + 0.32), Inches(8.9), Inches(0.55), a, size=14, color=MUTED)
        y += 1.0


def slide_backup_qa2(prs):
    slide = _blank_slide(prs)
    _title_bar(slide, "BACKUP · Q&A — limits / scale / data", "≤25s answers")

    pairs = [
        ("Limitations?", "Smoke fit is not_spec_floors · Verified promotion still future · landing % = preview"),
        ("Scale / business?", "Drop-in endpoint · per-hop FinOps log · teams + platforms · managed SaaS later"),
        ("Data / spend?", "2000 queries · ~1600/1000/800 gold/silver/cal · soft spend ≈ $8 · CI never spends"),
    ]
    y = 1.4
    for q, a in pairs:
        _text(slide, Inches(0.55), Inches(y), Inches(8.9), Inches(0.3), q, size=14, bold=True, color=ACCENT)
        _text(slide, Inches(0.55), Inches(y + 0.32), Inches(8.9), Inches(0.55), a, size=14, color=MUTED)
        y += 1.0


def slide_backup_clock(prs):
    slide = _blank_slide(prs)
    _title_bar(slide, "BACKUP · clock & judging", "Keep in your head")

    rows = [
        ("0:00–0:45", "Problem", "25%"),
        ("0:45–1:30", "Solution & innovation", "20%"),
        ("1:30–2:45", "Demo (one journey)", "15%"),
        ("2:45–3:25", "Technical", "20%"),
        ("3:25–4:00", "Potential & close", "20%"),
        ("+3:00", "Q&A", "—"),
    ]
    y = 1.4
    for clock, block, weight in rows:
        _text(slide, Inches(0.55), Inches(y), Inches(2.0), Inches(0.4), clock, size=14, bold=True, color=ACCENT)
        _text(slide, Inches(2.7), Inches(y), Inches(5.0), Inches(0.4), block, size=14)
        _text(slide, Inches(8.0), Inches(y), Inches(1.3), Inches(0.4), weight, size=14, color=MUTED)
        y += 0.5

    _text(
        slide,
        Inches(0.55),
        Inches(4.6),
        Inches(8.9),
        Inches(0.5),
        "If long: cut jargon first — never cut demo or closing.",
        size=14,
        color=OK,
    )


def main():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)

    slide_title(prs)
    slide_problem(prs)
    slide_solution(prs)
    slide_demo(prs)
    slide_algorithm(prs)
    slide_metrics(prs)
    slide_architecture(prs)
    slide_close(prs)
    slide_backup_qa1(prs)
    slide_backup_qa2(prs)
    slide_backup_clock(prs)

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    print(f"Wrote {OUT}")


if __name__ == "__main__":
    main()
