"""Generate hackathon submission PowerPoint from deck content."""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "lessons" / "0004-hackathon-submission.pptx"

# Pioneer-inspired palette
BG = RGBColor(0x0C, 0x06, 0x08)
INK = RGBColor(0xF5, 0xF5, 0xF4)
MUTED = RGBColor(0x8E, 0x8E, 0x96)
ACCENT = RGBColor(0xF2, 0x61, 0x3C)
BAD = RGBColor(0xFC, 0xA5, 0xA5)
OK = RGBColor(0x86, 0xEF, 0xAC)
CARD = RGBColor(0x16, 0x0A, 0x0E)
RULE = RGBColor(0x3F, 0x3F, 0x46)


def _blank_slide(prs: Presentation):
    layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(layout)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG
    return slide


def _box(slide, left, top, width, height, fill=CARD):
    shape = slide.shapes.add_shape(1, left, top, width, height)  # rectangle
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = RULE
    shape.line.width = Pt(0.75)
    return shape


def _text(slide, left, top, width, height, text, size=14, bold=False, color=INK, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = align
    return tf


def _bullets(tf, items, size=11, color=INK, bold_prefix=None):
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if bold_prefix and item.startswith(bold_prefix):
            p.text = item
        else:
            p.text = item
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.level = 0
        p.space_after = Pt(4)


def _title_bar(slide, title: str, subtitle: str = ""):
    _text(slide, Inches(0.55), Inches(0.35), Inches(9), Inches(0.55), title, size=28, bold=True)
    bar = slide.shapes.add_shape(1, Inches(0.55), Inches(0.95), Inches(1.2), Inches(0.06))
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()
    if subtitle:
        _text(slide, Inches(0.55), Inches(1.05), Inches(9), Inches(0.35), subtitle, size=12, color=MUTED)


def slide_cover(prs):
    slide = _blank_slide(prs)
    badges = "OpenAI-compatible  ·  Coding agents  ·  Cost-aware routing  ·  ML-trained path"
    _text(slide, Inches(0.8), Inches(1.6), Inches(8.4), Inches(0.4), badges, size=10, color=MUTED, align=PP_ALIGN.CENTER)
    _text(slide, Inches(0.8), Inches(2.2), Inches(8.4), Inches(0.9), "aiand/auto", size=44, bold=True, align=PP_ALIGN.CENTER)
    _text(
        slide,
        Inches(1.2),
        Inches(3.2),
        Inches(7.6),
        Inches(0.9),
        "A smart inference router that picks the cheapest aiand model\nper coding-agent step — without breaking hard tasks.",
        size=16,
        color=MUTED,
        align=PP_ALIGN.CENTER,
    )
    _text(slide, Inches(0.8), Inches(4.3), Inches(8.4), Inches(0.35), "Hackathon submission · Round One · 2026", size=11, color=MUTED, align=PP_ALIGN.CENTER)
    _text(
        slide,
        Inches(1.0),
        Inches(5.0),
        Inches(8.0),
        Inches(0.7),
        "One endpoint: model: router/auto\nDrop-in for OpenCode, Cursor, Claude Code, and any OpenAI client",
        size=13,
        align=PP_ALIGN.CENTER,
    )


def slide_problem(prs):
    slide = _blank_slide(prs)
    _title_bar(slide, "The Problem", "Problem fit & real-world impact")
    _text(
        slide,
        Inches(0.55),
        Inches(1.45),
        Inches(9),
        Inches(0.55),
        "AI coding agents run dozens of inference hops per session. Teams face a painful trade-off:",
        size=12,
    )

    left = _box(slide, Inches(0.55), Inches(2.05), Inches(4.35), Inches(1.65))
    left.line.color.rgb = BAD
    tf = _text(slide, Inches(0.7), Inches(2.15), Inches(4.05), Inches(1.5), "", size=11)
    tf.paragraphs[0].text = "Premium-only baseline"
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.size = Pt(12)
    _bullets(
        tf,
        [
            "Every step hits frontier models (Kimi K3, GLM-5)",
            "Simple steps cost 10–40× more than needed",
            "Budget caps stop agents mid-task",
        ],
        size=10,
    )

    right = _box(slide, Inches(5.1), Inches(2.05), Inches(4.35), Inches(1.65))
    right.line.color.rgb = BAD
    tf2 = _text(slide, Inches(5.25), Inches(2.15), Inches(4.05), Inches(1.5), "", size=11)
    tf2.paragraphs[0].text = "Cheap-only baseline"
    tf2.paragraphs[0].font.bold = True
    tf2.paragraphs[0].font.size = Pt(12)
    _bullets(
        tf2,
        [
            "Flash on everything fails multi-file refactors",
            "Tool-call accuracy drops on hard debug steps",
            "Human time lost re-running failed agent loops",
        ],
        size=10,
    )

    callout = _box(slide, Inches(0.55), Inches(3.85), Inches(8.9), Inches(0.75), CARD)
    callout.line.color.rgb = ACCENT
    _text(
        slide,
        Inches(0.7),
        Inches(3.95),
        Inches(8.6),
        Inches(0.65),
        "Who feels this? Engineering teams running OpenCode / Cursor agents against multi-model catalogs. "
        "Spend is measurable in dollars per session; quality in tests-passed and patch-applied.",
        size=10,
    )

    _text(slide, Inches(0.55), Inches(4.75), Inches(9), Inches(0.3), "Our solution in one sentence", size=12, bold=True)
    _text(
        slide,
        Inches(0.55),
        Inches(5.05),
        Inches(8.9),
        Inches(0.85),
        "aiand/auto sits between the agent and aiand as an OpenAI-compatible gateway. "
        "Each request carries phase context (x-agent-phase) and effort tier (x-routing-effort). "
        "The router filters ineligible models, scores survivors, and picks the cheapest model that clears your quality bar.",
        size=11,
        color=MUTED,
    )


def slide_built(prs):
    slide = _blank_slide(prs)
    _title_bar(slide, "What We Built", "Functional prototype — backend, frontend, ML pipeline")

    flow_y = Inches(1.45)
    steps = [
        ("Agent", "OpenCode / Cursor"),
        ("Gateway", "FastAPI /v1/chat"),
        ("Router", "Rules + shadow"),
        ("aiand", "9-model catalog"),
    ]
    x = 0.4
    for i, (title, sub) in enumerate(steps):
        bx = _box(slide, Inches(x), flow_y, Inches(2.15), Inches(0.7))
        _text(slide, Inches(x + 0.1), flow_y + Inches(0.05), Inches(1.95), Inches(0.25), title, size=11, bold=True)
        _text(slide, Inches(x + 0.1), flow_y + Inches(0.3), Inches(1.95), Inches(0.3), sub, size=9, color=MUTED)
        if i < len(steps) - 1:
            _text(slide, Inches(x + 2.2), flow_y + Inches(0.15), Inches(0.3), Inches(0.3), "→", size=14, color=ACCENT, align=PP_ALIGN.CENTER)
        x += 2.45

    _text(slide, Inches(0.55), Inches(2.35), Inches(4.3), Inches(0.25), "Backend (Python · FastAPI)", size=11, bold=True)
    tf = _text(slide, Inches(0.55), Inches(2.6), Inches(4.3), Inches(2.0), "", size=9)
    _bullets(
        tf,
        [
            "OpenAI-compatible proxy — streaming, tools, JSON",
            "Pioneer scoring formula on eligible models",
            "Hard constraints: tools, context, budget, latency",
            "Telemetry: data/requests.jsonl + /replay UI",
            "Eval harness vs premium / Kimi / adaptive",
            "70 automated tests — CI never spends credits",
        ],
        size=9,
        color=MUTED,
    )

    _text(slide, Inches(5.1), Inches(2.35), Inches(4.3), Inches(0.25), "Frontend (Next.js)", size=11, bold=True)
    tf2 = _text(slide, Inches(5.1), Inches(2.6), Inches(4.3), Inches(1.5), "", size=9)
    _bullets(
        tf2,
        [
            "Landing — simulator, ROI calculator",
            "Dashboard /routers/auto — savings, mix chart",
            "Playground — live router/auto chat",
            "Keys & usage — spend from request log",
        ],
        size=9,
        color=MUTED,
    )

    _text(slide, Inches(0.55), Inches(4.35), Inches(9), Inches(0.25), "ML pipeline (opt-in, budget-capped)", size=11, bold=True)
    pipeline = ["Teacher", "Gold", "Fit", "Shadow", "Trained"]
    subs = [
        "P(success) labels",
        "Measured trio",
        "Platt calibration",
        "Log would-pick",
        "Flip when sane",
    ]
    px = 0.55
    for label, sub in zip(pipeline, subs):
        _box(slide, Inches(px), Inches(4.65), Inches(1.65), Inches(0.55))
        _text(slide, Inches(px + 0.05), Inches(4.68), Inches(1.55), Inches(0.2), label, size=9, bold=True, align=PP_ALIGN.CENTER)
        _text(slide, Inches(px + 0.05), Inches(4.88), Inches(1.55), Inches(0.25), sub, size=7, color=MUTED, align=PP_ALIGN.CENTER)
        px += 1.75

    _text(
        slide,
        Inches(0.55),
        Inches(5.35),
        Inches(8.9),
        Inches(0.55),
        "Round One honesty: Live requests use the rules path. Trained routing runs in TRAINED_PATH=shadow — "
        "records X-Router-Trained-Would without changing the client response.",
        size=9,
        color=ACCENT,
    )


def slide_vs_rules(prs):
    slide = _blank_slide(prs)
    _title_bar(slide, "Why We Beat Rule-Based Routers", "Innovation & purposeful AI")

    left = _box(slide, Inches(0.55), Inches(1.45), Inches(4.35), Inches(2.35))
    left.line.color.rgb = BAD
    tf = _text(slide, Inches(0.7), Inches(1.55), Inches(4.05), Inches(2.2), "", size=9)
    tf.paragraphs[0].text = "Typical rule-based router"
    tf.paragraphs[0].font.bold = True
    _bullets(
        tf,
        [
            "Hand-tuned if-else: if phase=edit → Model X",
            "Static AA benchmarks, not your trajectories",
            "Same bar for typo fix and 12-file migration",
            "New model = retune dozens of thresholds",
            "Uncalibrated confidence — no real success rate",
        ],
        size=9,
        color=MUTED,
    )

    right = _box(slide, Inches(5.1), Inches(1.45), Inches(4.35), Inches(2.35))
    right.line.color.rgb = OK
    tf2 = _text(slide, Inches(5.25), Inches(1.55), Inches(4.05), Inches(2.2), "", size=9)
    tf2.paragraphs[0].text = "aiand/auto (rules → trained)"
    tf2.paragraphs[0].font.bold = True
    _bullets(
        tf2,
        [
            "Per-request P(success) from phase, tokens, tools",
            "Cheapest-above-bar + effort tiers + max-regret",
            "Complexity bins: trivial / standard / hard / frontier",
            "Platt calibration on gold completions (ECE ≤ 0.03)",
            "<10 ms in-process scorer — no LLM on hot path",
            "Graceful fallback: scorer_down → rules",
        ],
        size=9,
        color=MUTED,
    )

    _text(slide, Inches(0.55), Inches(3.95), Inches(9), Inches(0.25), "Three product pillars", size=11, bold=True)
    pillars = [
        ("01 · Score every request", "P(success) in <10 ms; no slow teacher on hot path"),
        ("02 · Clear confidence bar", "4 effort tiers; regret-bounded pool; dynamic pricing"),
        ("03 · One universal endpoint", "Drop-in /v1/chat/completions; X-Router-Model headers"),
    ]
    y = 4.25
    for title, desc in pillars:
        _text(slide, Inches(0.55), Inches(y), Inches(3.2), Inches(0.25), title, size=9, bold=True, color=ACCENT)
        _text(slide, Inches(3.5), Inches(y), Inches(5.9), Inches(0.25), desc, size=9, color=MUTED)
        y += 0.32

    callout = _box(slide, Inches(0.55), Inches(5.25), Inches(8.9), Inches(0.55))
    callout.line.color.rgb = OK
    _text(
        slide,
        Inches(0.7),
        Inches(5.32),
        Inches(8.6),
        Inches(0.45),
        "AI used purposefully: Teacher models label offline data; a tiny student classifier serves live traffic. "
        "No LLM call on every routing decision.",
        size=9,
    )


def slide_roadmap(prs):
    slide = _blank_slide(prs)
    _title_bar(slide, "Product Potential & Roadmap", "Scalability beyond the hackathon")

    _text(slide, Inches(0.55), Inches(1.45), Inches(9), Inches(0.25), "Target users", size=11, bold=True)
    tf = _text(slide, Inches(0.55), Inches(1.72), Inches(8.9), Inches(0.7), "", size=10)
    _bullets(
        tf,
        [
            "Platform teams — offer router/auto as managed endpoint on aiand",
            "Agent builders — OpenCode / Cursor shops wanting cost control",
            "FinOps / eng leads — dashboard proves savings with auditable JSONL",
        ],
        size=10,
        color=MUTED,
    )

    _text(slide, Inches(0.55), Inches(2.55), Inches(9), Inches(0.25), "Near-term (post-hackathon)", size=11, bold=True)
    timeline = [
        ("Now", "Rules live · shadow logging · dashboard + playground"),
        ("Next", "Smoke fit (teacher → gold → fit) · watch shadow vs rules"),
        ("Flip", "TRAINED_PATH=trained when shadow savings hold"),
        ("Gate", "SWE-bench Verified before production claim"),
    ]
    y = 2.85
    for when, what in timeline:
        _text(slide, Inches(0.55), Inches(y), Inches(0.9), Inches(0.25), when, size=9, bold=True, color=ACCENT)
        _text(slide, Inches(1.5), Inches(y), Inches(7.9), Inches(0.25), what, size=9, color=MUTED)
        y += 0.32

    _text(slide, Inches(0.55), Inches(4.15), Inches(4.3), Inches(0.25), "Commercialization", size=11, bold=True)
    tf_l = _text(slide, Inches(0.55), Inches(4.42), Inches(4.3), Inches(1.0), "", size=9)
    _bullets(
        tf_l,
        ["SaaS router with per-org retraining", "aiand partnership as default routing", "Usage-based pricing on measured savings"],
        size=9,
        color=MUTED,
    )

    _text(slide, Inches(5.1), Inches(4.15), Inches(4.3), Inches(0.25), "Technical scale", size=11, bold=True)
    tf_r = _text(slide, Inches(5.1), Inches(4.42), Inches(4.3), Inches(1.0), "", size=9)
    _bullets(
        tf_r,
        ["Multi-tenant gateway · production topology", "Retrain cadence on request cache", "Redaction + spend caps for enterprise"],
        size=9,
        color=MUTED,
    )

    _text(
        slide,
        Inches(0.55),
        Inches(5.45),
        Inches(8.9),
        Inches(0.35),
        "We never invent savings % — logged savings are vs most_expensive_eligible on each request.",
        size=9,
        color=MUTED,
    )


def slide_demo(prs):
    slide = _blank_slide(prs)
    _title_bar(slide, "Demo Guide & Evaluation Map", "How to present · how judges should read us")

    _text(slide, Inches(0.55), Inches(1.4), Inches(9), Inches(0.25), "5-minute live demo flow", size=11, bold=True)
    demo_rows = [
        ("1 · Hook", "Title / face", "Problem → one endpoint → aiand/auto"),
        ("2 · Landing", "localhost:3000", "Hero · 3 pillars · simulator"),
        ("3 · Dashboard", "/routers/auto", "Savings · mix chart · inference log"),
        ("4 · Playground", "/playground", "Live message · X-Router-Model"),
        ("5 · Future", "Architecture", "Shadow → trained · SWE-bench gate"),
    ]
    y = 1.68
    for step, screen, show in demo_rows:
        _text(slide, Inches(0.55), Inches(y), Inches(1.3), Inches(0.22), step, size=8, bold=True, color=ACCENT)
        _text(slide, Inches(1.9), Inches(y), Inches(1.8), Inches(0.22), screen, size=8, color=MUTED)
        _text(slide, Inches(3.8), Inches(y), Inches(5.6), Inches(0.22), show, size=8)
        y += 0.24

    _text(slide, Inches(0.55), Inches(2.95), Inches(9), Inches(0.25), "How we map to judging criteria", size=11, bold=True)
    criteria = [
        ("Problem fit & impact (25%)", "Agent cost/quality trade-off; measurable JSONL spend; real catalog prices"),
        ("Innovation & idea quality (20%)", "Student-teacher + cheapest-above-bar; AI offline not on hot path"),
        ("Technical implementation (20%)", "Gateway + dashboard + 70 tests; shadow/trained/scorer_down paths"),
        ("Product potential (20%)", "Drop-in endpoint; clear users; managed service roadmap"),
        ("Demo & presentation (15%)", "Landing + dashboard + playground in <3 min"),
    ]
    y = 3.22
    for label, detail in criteria:
        _text(slide, Inches(0.55), Inches(y), Inches(3.0), Inches(0.35), label, size=8, bold=True, color=ACCENT)
        _text(slide, Inches(3.6), Inches(y), Inches(5.8), Inches(0.35), detail, size=8, color=MUTED)
        y += 0.38

    _text(
        slide,
        Inches(0.55),
        Inches(5.2),
        Inches(8.9),
        Inches(0.45),
        "One-line pitch: “I built aiand/auto — an OpenAI-compatible router that picks the cheapest aiand model "
        "for each coding-agent step, with a landing page, live dashboard, and trained-ML path ready to flip on.”",
        size=10,
        color=INK,
    )


def main():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)  # 16:9

    slide_cover(prs)
    slide_problem(prs)
    slide_built(prs)
    slide_vs_rules(prs)
    slide_roadmap(prs)
    slide_demo(prs)

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    print(f"Wrote {OUT}")


if __name__ == "__main__":
    main()
